"""
File processing service for extracting text and handling various file formats.
"""

import base64
import io
from pathlib import Path
from pypdf import PdfReader
from docx import Document
from openpyxl import load_workbook
from pptx import Presentation
from PIL import Image


class FileProcessor:
    """Process various file types for AI consumption."""

    SUPPORTED_EXTENSIONS = {
        # Images
        '.jpg', '.jpeg', '.png', '.gif', '.bmp', '.webp',
        # Documents
        '.pdf', '.txt', '.md', '.docx',
        # Code
        '.py', '.js', '.ts', '.tsx', '.jsx', '.java', '.cpp', '.c', '.cs', '.go', '.rs',
        '.html', '.css', '.json', '.xml', '.yaml', '.yml', '.sh', '.sql',
        # Spreadsheets
        '.xlsx', '.csv',
        # Presentations
        '.pptx',
    }

    IMAGE_EXTENSIONS = {'.jpg', '.jpeg', '.png', '.gif', '.bmp', '.webp'}

    @classmethod
    def is_supported(cls, filename: str) -> bool:
        """Check if file type is supported."""
        ext = Path(filename).suffix.lower()
        return ext in cls.SUPPORTED_EXTENSIONS

    @classmethod
    def is_image(cls, filename: str) -> bool:
        """Check if file is an image."""
        ext = Path(filename).suffix.lower()
        return ext in cls.IMAGE_EXTENSIONS

    @classmethod
    def process_file(cls, file_content: bytes, filename: str) -> tuple[str, str | None]:
        """
        Process a file and extract text content.

        Returns:
            Tuple of (extracted_text, base64_image_data)
            - extracted_text: Text content extracted from the file
            - base64_image_data: Base64 encoded image data if file is an image, None otherwise
        """
        ext = Path(filename).suffix.lower()

        # Handle images
        if ext in cls.IMAGE_EXTENSIONS:
            return cls._process_image(file_content, ext)

        # Handle PDFs
        if ext == '.pdf':
            return cls._process_pdf(file_content), None

        # Handle Word documents
        if ext == '.docx':
            return cls._process_docx(file_content), None

        # Handle Excel files
        if ext == '.xlsx':
            return cls._process_xlsx(file_content), None

        # Handle CSV files
        if ext == '.csv':
            return file_content.decode('utf-8', errors='ignore'), None

        # Handle PowerPoint
        if ext == '.pptx':
            return cls._process_pptx(file_content), None

        # Handle plain text files
        if ext in {'.txt', '.md', '.py', '.js', '.ts', '.tsx', '.jsx', '.java',
                   '.cpp', '.c', '.cs', '.go', '.rs', '.html', '.css', '.json',
                   '.xml', '.yaml', '.yml', '.sh', '.sql'}:
            return file_content.decode('utf-8', errors='ignore'), None

        # Unsupported file type
        return f"[Unsupported file type: {ext}]", None

    @staticmethod
    def _process_image(file_content: bytes, ext: str) -> tuple[str, str | None]:
        """Process image file and return description + base64 data."""
        try:
            # Validate image
            img = Image.open(io.BytesIO(file_content))
            width, height = img.size
            mode = img.mode

            # Convert to RGB if necessary (for JPEG encoding)
            if mode not in ('RGB', 'RGBA'):
                img = img.convert('RGB')

            # Encode to base64
            buffer = io.BytesIO()
            img_format = 'PNG' if ext in {'.png', '.webp'} else 'JPEG'
            img.save(buffer, format=img_format)
            base64_data = base64.b64encode(buffer.getvalue()).decode('utf-8')

            # Create description
            description = f"[Image: {width}x{height}px, {mode} mode, {ext} format]"

            return description, base64_data

        except Exception as e:
            return f"[Error processing image: {str(e)}]", None

    @staticmethod
    def _process_pdf(file_content: bytes) -> str:
        """Extract text from PDF."""
        try:
            pdf_file = io.BytesIO(file_content)
            pdf_reader = PdfReader(pdf_file)

            text_parts = []
            for page_num, page in enumerate(pdf_reader.pages, 1):
                text = page.extract_text()
                if text.strip():
                    text_parts.append(f"--- Page {page_num} ---\n{text}")

            if not text_parts:
                return "[PDF contains no extractable text]"

            return "\n\n".join(text_parts)

        except Exception as e:
            return f"[Error processing PDF: {str(e)}]"

    @staticmethod
    def _process_docx(file_content: bytes) -> str:
        """Extract text from Word document."""
        try:
            doc_file = io.BytesIO(file_content)
            doc = Document(doc_file)

            text_parts = []

            # Extract paragraphs
            for para in doc.paragraphs:
                if para.text.strip():
                    text_parts.append(para.text)

            # Extract tables
            for table in doc.tables:
                table_text = []
                for row in table.rows:
                    row_text = ' | '.join(cell.text.strip() for cell in row.cells)
                    table_text.append(row_text)
                if table_text:
                    text_parts.append("\n".join(table_text))

            if not text_parts:
                return "[Document contains no text]"

            return "\n\n".join(text_parts)

        except Exception as e:
            return f"[Error processing DOCX: {str(e)}]"

    @staticmethod
    def _process_xlsx(file_content: bytes) -> str:
        """Extract text from Excel file."""
        try:
            xlsx_file = io.BytesIO(file_content)
            workbook = load_workbook(xlsx_file, data_only=True)

            text_parts = []

            for sheet_name in workbook.sheetnames:
                sheet = workbook[sheet_name]
                text_parts.append(f"=== Sheet: {sheet_name} ===\n")

                rows = []
                for row in sheet.iter_rows(values_only=True):
                    # Filter out empty rows
                    if any(cell is not None for cell in row):
                        row_text = ' | '.join(str(cell) if cell is not None else '' for cell in row)
                        rows.append(row_text)

                if rows:
                    text_parts.append("\n".join(rows))
                else:
                    text_parts.append("[Empty sheet]")

            if len(text_parts) == 0:
                return "[Workbook contains no data]"

            return "\n\n".join(text_parts)

        except Exception as e:
            return f"[Error processing XLSX: {str(e)}]"

    @staticmethod
    def _process_pptx(file_content: bytes) -> str:
        """Extract text from PowerPoint presentation."""
        try:
            pptx_file = io.BytesIO(file_content)
            presentation = Presentation(pptx_file)

            text_parts = []

            for slide_num, slide in enumerate(presentation.slides, 1):
                slide_text = [f"--- Slide {slide_num} ---"]

                for shape in slide.shapes:
                    if hasattr(shape, 'text') and shape.text.strip():
                        slide_text.append(shape.text)

                if len(slide_text) > 1:  # More than just the header
                    text_parts.append("\n".join(slide_text))

            if not text_parts:
                return "[Presentation contains no text]"

            return "\n\n".join(text_parts)

        except Exception as e:
            return f"[Error processing PPTX: {str(e)}]"
